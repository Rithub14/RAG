import os
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.document_loaders import Docx2txtLoader, TextLoader
from langchain_community.document_loaders.csv_loader import CSVLoader
from pptx import Presentation
import openpyxl
# from langchain_community.document_loaders import UnstructuredURLLoader
# from langchain_community.document_loaders import PlaywrightURLLoader
from langchain_community.document_loaders import WebBaseLoader



def load_file(file):
    name, extension = os.path.splitext(file)

    if extension == '.pdf':
        print("extension is ", extension)
        print(f'Loading {file}')
        loader = PyPDFLoader(file)

    elif extension == '.docx':
        print("extension is ", extension)
        print(f'Loading {file}')
        loader = Docx2txtLoader(file)

    elif extension == '.txt':
        print("extension is ", extension)
        print(f'Loading {file}')
        loader = TextLoader(file) 

    elif extension == '.csv':
        print("extension is ", extension)
        print(f'Loading {file}')
        loader = CSVLoader(file, autodetect_encoding=True)

    elif extension == '.xlsx':
        print("extension is ", extension)
        print(f'Loading {file}')
        workbook = openpyxl.load_workbook(filename=file, data_only=True)
        sheet = workbook.active
        all_text = ""
        for row in sheet.iter_rows():
            for cell in row:
                all_text += str(cell.value) + "\t"
            all_text += "\n"
        return all_text
    
    elif extension == '.pptx':
        print("extension is ", extension)
        print(f'Loading {file}')
        prs = Presentation(file)
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text + "\n"
        return all_text

    elif name[:5] == 'https':
        print("extension is ", name[:5])
        print(f'Loading {file}')
        # urls = [file]
        loader = WebBaseLoader(file)
        loader.requests_kwargs = {'verify':False}

    else:
        print('Format Not Supported.')
        return None

    data = loader.load()

    all_page_content = ""
    all_meta_data = ""
    for document in data:
        page_content = document.page_content
        all_page_content += page_content + "\n\n"

        metadata = document.metadata
        all_meta_data += str(metadata) + "\n\n"

    all_data = all_page_content + all_meta_data
    return all_data